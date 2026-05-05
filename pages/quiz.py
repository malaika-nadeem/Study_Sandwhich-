import streamlit as st
import pdfplumber
from docx import Document
from pptx import Presentation
import google.generativeai as genai
import io

st.set_page_config(page_title="studysandwhich", page_icon="🥪")
st.markdown("""
<style>
.stApp, [data-testid="stAppViewContainer"], [data-testid="stHeader"] {
    background-color: #FDF6EC !important;
}

.stButton button {
    background-color: #A0522D;
    color: white;
    border-radius: 20px;
    border: none;
}

@keyframes float {
    0% { transform: translateY(100vh) rotate(0deg); opacity: 0; }
    10% { opacity: 1; }
    90% { opacity: 1; }
    100% { transform: translateY(-100px) rotate(360deg); opacity: 0; }
}

.bubble {
    position: fixed;
    font-size: 30px;
    animation: float linear infinite;
    pointer-events: none;
    z-index: 0;
}

.b1 { left: 5%;  animation-duration: 6s;  animation-delay: 0s;  }
.b2 { left: 15%; animation-duration: 8s;  animation-delay: 1s;  }
.b3 { left: 30%; animation-duration: 7s;  animation-delay: 2s;  }
.b4 { left: 45%; animation-duration: 9s;  animation-delay: 0.5s;}
.b5 { left: 60%; animation-duration: 6s;  animation-delay: 3s;  }
.b6 { left: 75%; animation-duration: 8s;  animation-delay: 1.5s;}
.b7 { left: 88%; animation-duration: 7s;  animation-delay: 2.5s;}
</style>

<div class="bubble b1">🧋</div>
<div class="bubble b2">🥪</div>
<div class="bubble b3">🧋</div>
<div class="bubble b4">🥪</div>
<div class="bubble b5">🧋</div>
<div class="bubble b6">🥪</div>
<div class="bubble b7">🧋</div>
""", unsafe_allow_html=True)
st.markdown("""<style>
.stApp, [data-testid="stAppViewContainer"], [data-testid="stHeader"] {
    background-color: #FDF6EC !important;
.stButton button{
    background-color:#A0522D !important}           
}

[data-testid="stFileUploader"] {
    background-color: #A0522D;
    border: 2px solid #A0522D20;
    border-radius: 10px;
}

input[type="number"] {
    border: 2px solid ##A0522D20 !important;
    background-color: #A0522D !important;
    color: #3E2723 !important
                         

</style>""", unsafe_allow_html=True)
st.subheader("*Welcome to study room*🎓")
uploaded_file = st.file_uploader("Upload your notes", type=["pdf", "docx", "pptx", "txt"])
if uploaded_file is not None:
    st.markdown("""
        <div style="background-color:#A0522D20; border: 2px solid #A0522D; 
        border-radius:10px; padding:10px; color:#3E2723;">
        🥪 File uploaded successfully!
        </div>
    """, unsafe_allow_html=True)
    
    st.markdown("*Set the range of your notes*")
    page_start = st.number_input("Start-page", min_value=0, value=0)
    page_end = st.number_input("End-page", min_value=1, value=10)
    st.markdown("*Please select how many MCQ'S and ShortQuestion u want to generate🥪📖*")
    mcq=st.slider("Number of MCQs", 1, 10, 5)
    short=st.slider("Number of Short Question", 1, 10, 5)
    
    e_name = uploaded_file.name
    file_extension = e_name.split(".")[-1].lower()
    text = ""
    
    if file_extension == "txt":
        text = uploaded_file.read().decode("utf-8")
    
    elif file_extension == "pdf":
        with pdfplumber.open(io.BytesIO(uploaded_file.read())) as pdf:
            for page in pdf.pages[int(page_start):int(page_end)]:
                page_text = page.extract_text()
                if page_text:
                    text += page_text
    
    elif file_extension == "docx":
        doc = Document(io.BytesIO(uploaded_file.read()))
        text = "\n".join([p.text for p in doc.paragraphs])
    
    elif file_extension == "pptx":
        ppt = Presentation(io.BytesIO(uploaded_file.read()))
        for slide in ppt.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += shape.text + "\n"
    
    else:
        st.error("Unsupported file type!")

    if st.button("Start Quiz 🥪"):
     st.session_state["text"] = text
     st.session_state["num_mcq"] = mcq
     st.session_state["num_short"] = short
     st.switch_page("pages/test.py")
